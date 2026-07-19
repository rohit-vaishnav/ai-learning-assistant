from pptx import Presentation

def read_pptx(file_path: str) -> list[dict]:
    """
    Reads a PPTX file and extracts text slide-by-slide.
    Returns:
        list of dict: [{"text": slide_text, "page": slide_number}]
    """
    slides_data = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            text = "\n".join(slide_text)
            if text.strip():
                slides_data.append({
                    "text": text,
                    "page": i + 1
                })
    except Exception as e:
        raise RuntimeError(f"Error parsing PPTX: {str(e)}")
    return slides_data
